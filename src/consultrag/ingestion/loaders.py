"""Concrete loaders. Each is small and independently testable."""

from __future__ import annotations

from pathlib import Path

from .base import RawSection


class TextLoader:
    extensions = (".txt", ".md")

    def load(self, path: Path) -> list[RawSection]:
        text = path.read_text(encoding="utf-8", errors="ignore")
        return [RawSection(text=text, source_path=str(path), modality="text")]


class PptxLoader:
    """
    Extracts text from each slide, including titles, body placeholders, tables,
    and speaker notes. Speaker notes often hold the real analyst commentary, so
    we keep them and tag them.
    """
    extensions = (".pptx",)

    def load(self, path: Path) -> list[RawSection]:
        from pptx import Presentation

        prs = Presentation(str(path))
        sections: list[RawSection] = []
        for i, slide in enumerate(prs.slides, start=1):
            parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        parts.append(
                            " | ".join(cell.text for cell in row.cells)
                        )
            body = "\n".join(p for p in parts if p.strip())
            if body.strip():
                sections.append(
                    RawSection(
                        text=body,
                        source_path=str(path),
                        locator=f"slide {i}",
                        modality="slides",
                    )
                )
            # Speaker notes, if present.
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                if notes and notes.strip():
                    sections.append(
                        RawSection(
                            text=notes,
                            source_path=str(path),
                            locator=f"slide {i} (notes)",
                            modality="slides",
                            extra={"is_notes": True},
                        )
                    )
        return sections


class PdfLoader:
    """One section per page so we can cite page numbers."""
    extensions = (".pdf",)

    def load(self, path: Path) -> list[RawSection]:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        sections: list[RawSection] = []
        for i, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                sections.append(
                    RawSection(
                        text=text,
                        source_path=str(path),
                        locator=f"page {i}",
                        modality="pdf",
                    )
                )
        return sections


class AudioLoader:
    """
    Transcribes meeting recordings with faster-whisper, which runs fully locally
    (no audio ever leaves the machine — important for client confidentiality).

    Timestamps are preserved as the locator so an analyst can jump to the moment
    in the recording. Optional speaker diarization can be layered on with
    pyannote; left as an extension point to keep the default install light.

    Install: pip install faster-whisper
    """
    extensions = (".mp3", ".wav", ".m4a", ".flac", ".ogg")

    def __init__(self, model_size: str = "base", segment_group: int = 8):
        self.model_size = model_size
        self.segment_group = segment_group  # merge N segments per section

    def load(self, path: Path) -> list[RawSection]:
        from faster_whisper import WhisperModel  # local, CPU-friendly

        model = WhisperModel(self.model_size, device="cpu", compute_type="int8")
        segments, _info = model.transcribe(str(path))

        sections: list[RawSection] = []
        buffer: list[str] = []
        start_ts: float | None = None

        def _fmt(seconds: float) -> str:
            m, s = divmod(int(seconds), 60)
            h, m = divmod(m, 60)
            return f"{h:02d}:{m:02d}:{s:02d}"

        for idx, seg in enumerate(segments):
            if start_ts is None:
                start_ts = seg.start
            buffer.append(seg.text.strip())
            if (idx + 1) % self.segment_group == 0:
                sections.append(
                    RawSection(
                        text=" ".join(buffer),
                        source_path=str(path),
                        locator=_fmt(start_ts),
                        modality="audio",
                    )
                )
                buffer, start_ts = [], None
        if buffer:
            sections.append(
                RawSection(
                    text=" ".join(buffer),
                    source_path=str(path),
                    locator=_fmt(start_ts or 0),
                    modality="audio",
                )
            )
        return sections
